# -*- coding: utf-8 -*-
"""
Created on Sun Apr 29 20:54:18 2018

@author: Administrator
"""

# 测试框架
import unittest

# PPT
from pptx import Presentation

import auto_report.hello_pptx as hello_pptx 

class TestHelloPptx(unittest.TestCase):

    def test_add_all_slide_layout(self):
        prs = Presentation('../resources/template/ppt_template0.pptx')
        hello_pptx.add_all_slide_layout(prs)
        prs.save('../resources/report/slide_layout.pptx')
    
    def test_add_shape_text(self):
        prs = Presentation('../resources/template/ppt_template0.pptx')
        
        
        slide_layout = prs.slide_layouts[4]
        
        # add_text_by_shape
        slide = prs.slides.add_slide(slide_layout)
        hello_pptx.add_text_by_shape(slide)
        
        # add_text_by_palceholder
        slide = prs.slides.add_slide(slide_layout)
        hello_pptx.add_text_by_palceholder(slide)
        
        
        prs.save('../resources/report/add_shape_text.pptx')
    
            