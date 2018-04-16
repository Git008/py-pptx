# -*- coding: utf-8 -*-
"""
Created on Thu Apr  5 12:12:34 2018

@author: Administrator
"""

# 数据分析
import matplotlib.pyplot as plt
import pandas as pd

# PPT
from pptx import Presentation
from pptx.util import Inches 

# web 
from selenium import webdriver

# excel
import xlwt
import time

# slide layout测试
def test_slide_layout(prs):
    slide_layout_list = ['Title (presentation title slide)',
                         'Title and Content',
                         'Section Header (sometimes called Segue)',
                         'Two Content (side by side bullet textboxes)',
                         'Comparison (same but additional title for each side by side content box)',
                         'Title Only',
                         'Blank',
                         'Content with Caption',
                         'Picture with Caption']
    
    for index, name in enumerate(slide_layout_list):
        new_slide_layout = prs.slide_layouts[index]
        new_slide = prs.slides.add_slide(new_slide_layout)
        
        #6对应的是Blank，没有标题框
        if index != 6:
            title = new_slide.shapes.title
            title.text = name

# 箱线图    
def page1_boxplot(prs):
    # slide_layouts[1]为带标题和正文框的ppt
    title_slide_layout = prs.slide_layouts[1]
    
    # 新增ppt（当前仅支持新增操作）
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 取本页ppt的title，向title文本框写如文字 
    title = slide.shapes.title
    title.text = '箱线图示例' 
    
    # 导入数据集
    df = pd.read_excel('./resource/data0.xlsx', 'Sheet1')
    
    # 绘制箱线图 
    # 方法1：使用matplotlib，暂未找到方法，待完善
    # fig = plt.figure()
    # ax = fig.add_subplot(111)
    # ax.boxplot(df['Lead Time'])
    
    # 方法2：使用dataframe，第一个参数column是数据列名，第二个参数by是分组的列名
    df.boxplot(column='Lead Time', by='Project')
    
    plt.savefig('./resource./交付周期_箱线图.png') 
    plt.show()
    
    # 插入图片，在指定位置按预设值添加图片 
    img_path = './resource./交付周期_箱线图.png' 
    left, top, width, height = Inches(1), Inches(1.8), Inches(4.5), Inches(4.5)  
    slide.shapes.add_picture(img_path, left, top, width, height)

def hello_selenium():
    # chromedriverv2.3.7 配套chromev65
    browser = webdriver.Chrome('./driver/chromedriver2.3.7.exe')
    
    # 打开指定网页
    url = "http://weixin.sogou.com/"  
    browser.get(url)
    
    # 在微信搜索框中输入关键字"六西格玛是个P"，通过ID查找
    browser.find_element_by_id('query').send_keys('六西格玛是个P')
    # 单击公众号搜索按钮，通过class查找
    browser.find_element_by_class_name('swz2').click()
    
    # 关闭浏览器
    browser.quit()
    
def hello_table_data():
    # chromedriverv2.3.7 配套chromev65
    browser = webdriver.Chrome('./driver/chromedriver2.3.7.exe')
    
    # 打开指定网
    browser.get('http://newhouse.nj.house365.com/')
    
    # 等待5秒，页面加载OK，解决NoSuchElementException问题
    time.sleep(5)    
    
    #创建工作簿  
    wbk = xlwt.Workbook(encoding='utf-8', style_compression=0)  
    #创建工作表  
    sheet = wbk.add_sheet('sheet 1', cell_overwrite_ok=True)  
 
    #今日交易信息快递
    #将表的每一行存在table_tr_list中
    table_body_xpath = '/html/body/div[19]/div[2]/div[2]/div[1]/table/tbody'
    table_body = browser.find_element_by_xpath(table_body_xpath)
    table_tr_list = table_body.find_elements_by_tag_name('tr')  
    
    # 从excel第一行开始存
    for r,tr in enumerate(table_tr_list, 0):  
        #将表的每一行的每一列内容存在table_td_list中  
        table_td_list = tr.find_elements_by_tag_name('td')  
        #写入表的内容到sheet 1中，第r行第c列  
        for c,td in enumerate(table_td_list):  
            sheet.write(r, c, td.text) 
            
    #保存表格到已有的 excel  
    wbk.save('test.xls')  
    
    # 关闭浏览器
    browser.quit() 

if __name__ == "__main__":
   # 使用自定义模板
    prs = Presentation('./template/ppt_template0.pptx')
    
    # page1_boxplot(prs)
    hello_table_data()
    
    # 保存ppt
    prs.save('数据分析报告.pptx')